# backend/extractor.py
"""Universal file text extractor. Never raises -- always returns str."""
import os
import subprocess
from pathlib import Path
from typing import Optional


def extract_text(file_path: str) -> str:
    """Extract text from any file. Returns empty string if extraction fails."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _extract_pdf(file_path)
        elif suffix in (".docx", ".doc"):
            return _extract_docx(file_path)
        elif suffix in (".txt", ".md", ".csv"):
            return path.read_text(errors="replace")
        elif suffix == ".rtf":
            return _extract_rtf(file_path)
        elif suffix in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif suffix in (".xlsx", ".xls"):
            return _extract_xlsx(file_path)
        elif suffix in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp"):
            return _extract_image(file_path)
        else:
            # Try UTF-8 read, fall back to Vision
            try:
                text = path.read_text(errors="strict")
                if len(text.strip()) > 20:
                    return text
            except Exception:
                pass
            return _extract_image(file_path)
    except Exception:
        return ""


def _extract_pdf(file_path: str) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            text_parts.append(text)
    full_text = "\n".join(text_parts).strip()
    if len(full_text) < 50 * max(1, len(text_parts)):
        return _extract_image(file_path)
    return full_text


def _extract_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_rtf(file_path: str) -> str:
    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", file_path],
        capture_output=True, text=True, timeout=30
    )
    return result.stdout if result.returncode == 0 else ""


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_image(file_path: str) -> str:
    """Use Claude Vision to extract text from image or scanned document."""
    import anthropic
    import base64
    path = Path(file_path)
    suffix = path.suffix.lower()
    media_map = {
        ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".png": "image/png", ".tiff": "image/tiff", ".tif": "image/tiff",
        ".bmp": "image/bmp", ".webp": "image/webp", ".pdf": "application/pdf",
    }
    media_type = media_map.get(suffix, "image/png")
    with open(file_path, "rb") as f:
        data = base64.standard_b64encode(f.read()).decode()
    client = anthropic.Anthropic()
    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [{
                "type": "image",
                "source": {"type": "base64", "media_type": media_type, "data": data}
            }, {
                "type": "text",
                "text": "Extract all readable text from this document. Return only the text, preserving section structure."
            }]
        }]
    )
    return response.content[0].text
